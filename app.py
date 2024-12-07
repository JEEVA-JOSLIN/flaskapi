# from flask import Flask,request,jsonify
# import json
# import base64
# import magic

# app = Flask(__name__)

# @app.route('/')
# def home():
#     return "Hello from Flask API on Azure!"

# @app.route('/process', methods=['POST'])

# def process_file():
#     if 'file' not in request.files:
#         return jsonify({"error": "No file provided"}), 400
#     file = request.files['file']
#     if file.filename == '':
#         return jsonify({"error": "No selected file"}), 400

#     # Save the uploaded file temporarily
#     file_path = f"./{file.filename}"
#     file.save(file_path)

#     # Using the magic module to identify the file type
#     file_format = identify_file_format(file_path)
    
#     # Creating a sample response with file type and base64 encoding
#     content = {
#         "file_name": file.filename,
#         "file_format": file_format,
#         "base64_content": encode_file_to_base64(file_path)
#     }
    
#     return jsonify(content)

# # Function to identify file type using magic module
# def identify_file_format(file_path):
#     mime = magic.Magic(mime=True)
#     file_format = mime.from_file(file_path)
#     return file_format

# # Function to convert file to base64 encoding
# def encode_file_to_base64(file_path):
#     with open(file_path, "rb") as file:
#         encoded_string = base64.b64encode(file.read()).decode("utf-8")
#     return encoded_string

# if __name__ == "__main__":
#     app.run(debug=True)

from flask import Flask, request, jsonify
import json
import base64
import magic
import os

# Flask setup
app = Flask(__name__)

# Directory to save uploaded files temporarily
UPLOAD_FOLDER = './uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# MasterApp class encapsulates file processing logic
class MasterApp:
    def __init__(self):
        self.magic = magic.Magic()

    def process_file(self, file_path):
        file_format = self.identify_file_format(file_path)
        content = {}
        if "text" in file_format.lower():
            content = self.process_text(file_path)
        elif "pdf" in file_format.lower():
            content = self.process_pdf(file_path)
        elif "word" in file_format.lower():
            content = self.process_docx(file_path)
        elif "powerpoint" in file_format.lower():
            content = self.process_pptx(file_path)
        else:
            return {"error": "Unsupported file type."}
        output_path = os.path.splitext(file_path)[0] + ".json"
        self.save_to_json(content, output_path)
        return content

    def identify_file_format(self, file_path):
        return self.magic.from_file(file_path)

    def save_to_json(self, content, output_path):
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=4)

    def process_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            text_content = file.read()
            return {
                "page_1": {
                    "text": {"content": text_content},
                    "recognized_text": {},
                    "images": [],
                    "tables": {}
                }
            }

    def process_pdf(self, file_path):
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        content = {}
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            images = []
            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                image_base64 = base64.b64encode(img_bytes).decode('utf-8')
                images.append({"base64": image_base64})
            content[f"page_{page_num + 1}"] = {
                "text": {"content": text},
                "recognized_text": {},
                "images": images,
                "tables": {}
            }
        doc.close()
        return content

    def process_docx(self, file_path):
        from docx import Document
        doc = Document(file_path)
        content = {}
        page_counter = 1
        for para in doc.paragraphs:
            content[f"page_{page_counter}"] = {
                "text": {"content": para.text.strip()},
                "recognized_text": {},
                "images": [],
                "tables": {}
            }
            page_counter += 1
        return content

    def process_pptx(self, file_path):
        from pptx import Presentation
        presentation = Presentation(file_path)
        content = {}
        page_counter = 1
        for slide in presentation.slides:
            text_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text_content.append(para.text.strip())
            content[f"page_{page_counter}"] = {
                "text": {"content": " ".join(text_content)},
                "recognized_text": {},
                "images": [],
                "tables": {}
            }
            page_counter += 1
        return content

# Initialize the MasterApp
app.master_app = MasterApp()

@app.route('/')
def home():
    return "Hello from Flask API on Azure!"

@app.route('/process', methods=['POST'])
def process_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    # Save the uploaded file temporarily
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
    file.save(file_path)

    # Process the file using the MasterApp class
    content = app.master_app.process_file(file_path)
    
    # Return the content in JSON format
    return jsonify(content)

if __name__ == "__main__":
    app.run(debug=True)

